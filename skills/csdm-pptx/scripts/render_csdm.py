#!/usr/bin/env python3
"""
render_csdm.py - Deterministic CSDM v5 slide renderer.

Reads a spec JSON describing which CSDM objects, roles and relationships to
visualise and produces a standalone .pptx with one slide that places every
requested object at its canonical CSDM coordinate with the correct domain
colour, then draws the relationships as labelled connectors.

Usage:
    python render_csdm.py <spec.json> <output.pptx>
    python render_csdm.py --spec <spec.json> --out <output.pptx>
                          [--catalog <csdm_model.json>]
                          [--hide-empty-domains]
                          [--dump-catalog]  # prints known objects/roles and exits

Spec schema:
{
  "title":   "My CSDM view",
  "subtitle":"Optional second line",
  "objects": ["Business Application", "Information Object", ...],
  "roles":   ["Digital Product Owner", ...],             # optional
  "relationships": [
      {"from": "Service Instance", "to": "Business Application",
       "label": "hosts", "style": "straight"|"bent"|"curved"}
  ],
  "hide_empty_domains": false,   # if true, domains with no mentioned object are hidden
  "show_domain_labels": true,
  "custom_objects": [            # optional: objects not in the canonical catalog
      {"name": "Custom Thing", "domain": "service_delivery",
       "position": {"x": 3.0, "y": 3.0, "w": 1.0, "h": 0.4},
       "fill": "FFFFFF"}
  ],
  "domain_overrides": {          # optional: resize/move domain regions + labels
      "design_planning": {
          "region": {"x": 7.0, "y": 0.8, "w": 6.0, "h": 2.2,
                     "shape": "roundRect"},
          "label_position": {"x": 7.0, "y": 0.5, "w": 2.0, "h": 0.30,
                             "anchor": "outside_left"}
      }
  }
}

Object / role names are matched case-insensitively against the canonical
name and all synonyms. If a name cannot be matched the script fails with
a list of near matches.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError as exc:
    sys.stderr.write(
        "ERROR: python-pptx is required. Install with:\n"
        "  pip install python-pptx --break-system-packages\n"
    )
    raise

DEFAULT_CATALOG = Path(__file__).parent.parent / "assets" / "csdm_model.json"


# --------------------------------------------------------------------------- #
# Catalog helpers
# --------------------------------------------------------------------------- #
def load_catalog(path: Path) -> dict:
    with open(path) as f:
        return json.load(f)


def _normalise(s: str) -> str:
    return re.sub(r"[^a-z0-9]", "", s.lower())


def build_lookup(catalog: dict) -> Tuple[Dict[str, dict], Dict[str, dict]]:
    """Build normalised-name → entry lookups for objects and roles."""
    obj_lookup: Dict[str, dict] = {}
    for o in catalog["objects"]:
        obj_lookup[_normalise(o["canonical"])] = o
        for syn in o.get("synonyms", []):
            obj_lookup[_normalise(syn)] = o
    role_lookup: Dict[str, dict] = {}
    for r in catalog["roles"]:
        role_lookup[_normalise(r["canonical"])] = r
        for syn in r.get("synonyms", []):
            role_lookup[_normalise(syn)] = r
    return obj_lookup, role_lookup


def resolve(name: str, lookup: Dict[str, dict]) -> Optional[dict]:
    return lookup.get(_normalise(name))


def domain_by_id(catalog: dict, did: str) -> Optional[dict]:
    for d in catalog["domains"]:
        if d["id"] == did:
            return d
    return None


# --------------------------------------------------------------------------- #
# Rendering primitives
# --------------------------------------------------------------------------- #
SHAPE_MAP = {
    "rect":         MSO_SHAPE.RECTANGLE,
    "roundRect":    MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse":      MSO_SHAPE.OVAL,
    "chevron":      MSO_SHAPE.CHEVRON,
    "cloudCallout": MSO_SHAPE.CLOUD_CALLOUT,
}


def rgb(hex_: str) -> RGBColor:
    return RGBColor(int(hex_[0:2], 16), int(hex_[2:4], 16), int(hex_[4:6], 16))


def add_shape(slide, kind: str, x, y, w, h, fill_hex=None, stroke_hex=None,
              text=None, font_size=9, font_color="000000", bold=False,
              align_center=True):
    mso = SHAPE_MAP.get(kind, MSO_SHAPE.RECTANGLE)
    shp = slide.shapes.add_shape(mso, Inches(x), Inches(y), Inches(w), Inches(h))

    if fill_hex is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill_hex)

    if stroke_hex:
        shp.line.color.rgb = rgb(stroke_hex)
        shp.line.width = Pt(0.75)
    else:
        shp.line.color.rgb = rgb("808080")
        shp.line.width = Pt(0.5)

    tf = shp.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.word_wrap = True

    if text is None:
        text = ""
    tf.text = text
    for p in tf.paragraphs:
        if align_center:
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.color.rgb = rgb(font_color)
            r.font.name = "Calibri"
    # Vertical center
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shp


def add_textbox(slide, x, y, w, h, text, font_size=12, bold=False,
                color="000000", align="center"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.text = text
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    for p in tf.paragraphs:
        p.alignment = {"center": PP_ALIGN.CENTER,
                       "left":   PP_ALIGN.LEFT,
                       "right":  PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.color.rgb = rgb(color)
            r.font.name = "Calibri"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


# --------------------------------------------------------------------------- #
# Connector rendering
# --------------------------------------------------------------------------- #
def _shape_anchors(shape) -> Dict[str, Tuple[int, int]]:
    """Return EMU coordinates for L, R, T, B, C anchors of a shape."""
    left, top = shape.left, shape.top
    w, h = shape.width, shape.height
    return {
        "L": (left,         top + h // 2),
        "R": (left + w,     top + h // 2),
        "T": (left + w // 2, top),
        "B": (left + w // 2, top + h),
        "C": (left + w // 2, top + h // 2),
    }


def _pick_sides(from_shape, to_shape) -> Tuple[str, str]:
    fa = _shape_anchors(from_shape)
    ta = _shape_anchors(to_shape)
    fc = fa["C"]
    tc = ta["C"]
    dx = tc[0] - fc[0]
    dy = tc[1] - fc[1]
    if abs(dx) >= abs(dy):
        return ("R", "L") if dx >= 0 else ("L", "R")
    else:
        return ("B", "T") if dy >= 0 else ("T", "B")


def add_connector(slide, from_shape, to_shape, style="straight",
                  label: Optional[str] = None, color="000000",
                  width_pt=1.25, with_arrow=True):
    fa = _shape_anchors(from_shape)
    ta = _shape_anchors(to_shape)
    fs, ts = _pick_sides(from_shape, to_shape)
    x1, y1 = fa[fs]
    x2, y2 = ta[ts]

    if style == "bent":
        conn_type = MSO_CONNECTOR.ELBOW
    elif style == "curved":
        conn_type = MSO_CONNECTOR.CURVE
    else:
        conn_type = MSO_CONNECTOR.STRAIGHT

    conn = slide.shapes.add_connector(conn_type, x1, y1, x2, y2)
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width_pt)

    # Bind endpoints to the shapes so connectors follow the shapes and render
    # as properly-anchored lines (no free-floating "green dot" endpoints in
    # LibreOffice / PowerPoint). Connection site indices for a rectangle are:
    #   0 = top, 1 = right, 2 = bottom, 3 = left
    # (python-pptx uses the same convention as OOXML cxnLst ordering for
    # standard auto-shapes.)
    SIDE_TO_SITE = {"T": 0, "R": 1, "B": 2, "L": 3}
    try:
        conn.begin_connect(from_shape, SIDE_TO_SITE.get(fs, 1))
        conn.end_connect(to_shape, SIDE_TO_SITE.get(ts, 3))
    except (AttributeError, KeyError, ValueError):
        # If the python-pptx version does not expose begin_connect/end_connect
        # or the target shape has no connection sites, fall through with the
        # absolute-coordinate connector we already created.
        pass

    # Arrow head
    if with_arrow:
        ln = conn.line._get_or_add_ln()
        tail_end = etree.SubElement(
            ln, qn("a:tailEnd"),
            {"type": "triangle", "w": "med", "len": "med"}
        )

    # Label - place a transparent textbox midway
    if label:
        mx_in = (x1 + x2) / 2 / 914400
        my_in = (y1 + y2) / 2 / 914400
        lbl_w = max(0.8, min(1.6, len(label) * 0.08))
        lbl_h = 0.22
        tb = slide.shapes.add_textbox(
            Inches(mx_in - lbl_w / 2), Inches(my_in - lbl_h / 2),
            Inches(lbl_w), Inches(lbl_h)
        )
        tf = tb.text_frame
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        tf.text = label
        from pptx.enum.text import PP_ALIGN
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(8)
                r.font.italic = True
                r.font.color.rgb = rgb("404040")
                r.font.name = "Calibri"
    return conn


# --------------------------------------------------------------------------- #
# Role icon (circle with initials)
# --------------------------------------------------------------------------- #
def add_role(slide, role_entry: dict):
    pos = role_entry["position"]
    x, y, w, h = pos["x"], pos["y"], pos["w"], pos["h"]
    # Circle with person silhouette look - use filled ellipse with initials
    icon_size = min(h, 0.42)
    # Place icon left of the label box
    ic_x = x - icon_size - 0.05
    if role_entry.get("anchor") == "left":
        ic_x = x + w + 0.05
    elif role_entry.get("anchor") == "below":
        ic_x = x + (w - icon_size) / 2
        y -= icon_size + 0.02
    ic = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(ic_x), Inches(y),
                                Inches(icon_size), Inches(icon_size))
    ic.fill.solid()
    ic.fill.fore_color.rgb = rgb("5B9BD5")
    ic.line.color.rgb = rgb("2E75B6")
    ic.line.width = Pt(0.5)
    # Initials
    initials = "".join(part[0] for part in role_entry["canonical"].split() if part[:1].isalpha())[:3]
    tf = ic.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.text = initials
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(8)
            r.font.bold = True
            r.font.color.rgb = rgb("FFFFFF")
            r.font.name = "Calibri"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Role label
    add_textbox(slide, x, y, w, h,
                role_entry["canonical"].replace(" (Foundation)", "")
                                       .replace(" (role)", ""),
                font_size=8, bold=False, color="000000",
                align="left" if role_entry.get("anchor") == "right" else "center")


# --------------------------------------------------------------------------- #
# Main render
# --------------------------------------------------------------------------- #
def render(spec: dict, catalog: dict, out_path: Path) -> dict:
    obj_lookup, role_lookup = build_lookup(catalog)

    # Resolve requested objects
    resolved_objects: List[dict] = []
    unknown: List[str] = []
    for name in spec.get("objects", []):
        entry = resolve(name, obj_lookup)
        if entry is None:
            unknown.append(name)
        else:
            if entry not in resolved_objects:
                resolved_objects.append(entry)

    # Custom objects pass through
    for co in spec.get("custom_objects", []) or []:
        resolved_objects.append({
            "canonical": co["name"],
            "domain": co.get("domain"),
            "position": co["position"],
            "shape": co.get("shape", "rect"),
            "fill": co.get("fill", "FFFFFF"),
            "_custom": True,
        })

    # Resolve requested roles
    resolved_roles: List[dict] = []
    for name in spec.get("roles", []) or []:
        entry = resolve(name, role_lookup)
        if entry is None:
            unknown.append(f"role:{name}")
        else:
            if entry not in resolved_roles:
                resolved_roles.append(entry)

    if unknown:
        known_obj = sorted({o["canonical"] for o in catalog["objects"]})
        known_role = sorted({r["canonical"] for r in catalog["roles"]})
        msg = (
            "Unknown CSDM names:\n  "
            + "\n  ".join(unknown)
            + "\n\nKnown objects:\n  "
            + "\n  ".join(known_obj)
            + "\n\nKnown roles:\n  "
            + "\n  ".join(known_role)
        )
        raise ValueError(msg)

    # Build presentation
    slide_w = catalog["slide"]["width_in"]
    slide_h = catalog["slide"]["height_in"]
    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)
    blank_layout = prs.slide_layouts[6]  # usually "Blank"
    slide = prs.slides.add_slide(blank_layout)

    # Title (positioned in the top-left corner, left of Build & Integration region which starts at x=4.39)
    title = spec.get("title") or "CSDM v5 - Core Data Objects"
    subtitle = spec.get("subtitle") or ""
    add_textbox(slide, 0.25, 0.10, 4.00, 0.35, title,
                font_size=16, bold=True, color="203864", align="left")
    if subtitle:
        add_textbox(slide, 0.25, 0.45, 4.00, 0.25, subtitle,
                    font_size=10, bold=False, color="595959", align="left")

    hide_empty = bool(spec.get("hide_empty_domains", False))

    # Which domains actually contain mentioned objects
    used_domains = {o.get("domain") for o in resolved_objects if o.get("domain")}
    # Always show manage_portfolio and foundation if any object is shown
    if resolved_objects:
        used_domains.add("manage_portfolio")
        used_domains.add("foundation")

    # Optional per-domain overrides supplied in the spec. Allows callers to
    # resize/reposition domain regions and their labels when the default
    # canonical layout is too cramped. Format:
    #   "domain_overrides": {
    #       "design_planning": {
    #           "region": {"x": 7.0, "y": 0.8, "w": 6.0, "h": 2.2},
    #           "label_position": {"x": 7.0, "y": 0.5, "w": 2.0, "h": 0.30,
    #                              "anchor": "outside_left"}
    #       }
    #   }
    domain_overrides = spec.get("domain_overrides", {}) or {}

    def _region_for(d):
        ov = domain_overrides.get(d["id"], {})
        return ov.get("region", d["region"])

    def _label_for(d):
        ov = domain_overrides.get(d["id"], {})
        return ov.get("label_position", d.get("label_position"))

    # 1) Draw domain regions (back layer)
    for d in catalog["domains"]:
        if hide_empty and d["id"] not in used_domains and d["id"] != "manage_portfolio":
            continue
        r = _region_for(d)
        shp = slide.shapes.add_shape(
            SHAPE_MAP.get(r["shape"], MSO_SHAPE.ROUNDED_RECTANGLE),
            Inches(r["x"]), Inches(r["y"]), Inches(r["w"]), Inches(r["h"])
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(d["fill"])
        stroke = d.get("stroke", d["fill"])
        shp.line.color.rgb = rgb(stroke)
        shp.line.width = Pt(0.75)
        shp.text_frame.text = ""  # no text in backdrop

    # 2) Domain labels (above backdrop). Labels anchored "outside_*" get a
    # white background so they read cleanly when they happen to sit near a
    # coloured region edge. "inside" labels (Manage Portfolio, Foundation)
    # have no background so they blend into the host shape.
    if spec.get("show_domain_labels", True):
        for d in catalog["domains"]:
            if hide_empty and d["id"] not in used_domains and d["id"] != "manage_portfolio":
                continue
            lp = _label_for(d)
            if not lp or not d["label"]:
                continue
            anchor = lp.get("anchor", "outside_left")
            # Draw a white rounded pill behind outside labels so they read
            # cleanly against any background
            if anchor.startswith("outside"):
                bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(lp["x"] - 0.05), Inches(lp["y"] - 0.03),
                    Inches(lp["w"] + 0.10), Inches(lp["h"] + 0.06)
                )
                bg.fill.solid()
                bg.fill.fore_color.rgb = rgb("FFFFFF")
                bg.line.color.rgb = rgb(d.get("stroke", d["fill"]))
                bg.line.width = Pt(0.75)
                bg.text_frame.text = ""
            add_textbox(slide, lp["x"], lp["y"], lp["w"], lp["h"],
                        d["label"], font_size=10, bold=True,
                        color=d.get("text_color", "000000"),
                        align="left" if anchor == "inside" and d["id"] == "foundation" else "center")

    # 3) Draw objects
    name_to_shape: Dict[str, object] = {}
    for o in resolved_objects:
        pos = o["position"]
        display = o["canonical"]
        # Strip parenthetical disambiguators from display text
        display = re.sub(r"\s*\(lower\)\s*$", "", display)
        # Stacked objects: draw N rectangles slightly offset
        stacked = int(o.get("stacked", 1))
        fill = o.get("fill", "FFFFFF")
        kind = o.get("shape", "rect")
        x, y, w, h = pos["x"], pos["y"], pos["w"], pos["h"]
        last_shape = None
        for s in range(stacked):
            off = s * 0.08
            sh = add_shape(slide, kind, x + off, y + off, w, h,
                           fill_hex=fill, stroke_hex="404040",
                           text=(display if s == stacked - 1 else ""),
                           font_size=9, font_color="000000", bold=False)
            last_shape = sh
        name_to_shape[_normalise(o["canonical"])] = last_shape

    # 4) Draw roles
    for r in resolved_roles:
        add_role(slide, r)

    # 5) Draw relationships
    rels = spec.get("relationships", []) or []
    # If no explicit relationships, optionally include standard ones between
    # drawn objects
    if not rels and spec.get("auto_relationships", True):
        drawn_names = {_normalise(o["canonical"]) for o in resolved_objects}
        for rel in catalog.get("standard_relationships", []):
            if _normalise(rel["from"]) in drawn_names and _normalise(rel["to"]) in drawn_names:
                rels.append(rel)

    rel_errors: List[str] = []
    for rel in rels:
        f_key = _normalise(rel["from"])
        t_key = _normalise(rel["to"])
        # Resolve names via catalog in case user used a synonym
        f_entry = resolve(rel["from"], obj_lookup)
        t_entry = resolve(rel["to"], obj_lookup)
        if f_entry is not None:
            f_key = _normalise(f_entry["canonical"])
        if t_entry is not None:
            t_key = _normalise(t_entry["canonical"])
        fs = name_to_shape.get(f_key)
        ts = name_to_shape.get(t_key)
        if fs is None or ts is None:
            rel_errors.append(
                f"Skipping relationship {rel['from']} -> {rel['to']} "
                f"(one side not drawn)"
            )
            continue
        style = rel.get("style")
        if not style:
            # map catalog 'connector' field
            conn = rel.get("connector", "straightConnector1")
            style = "bent" if "bent" in conn else ("curved" if "curved" in conn else "straight")
        add_connector(slide, fs, ts, style=style, label=rel.get("label"))

    prs.save(out_path)

    return {
        "output": str(out_path),
        "objects_drawn": [o["canonical"] for o in resolved_objects],
        "roles_drawn": [r["canonical"] for r in resolved_roles],
        "relationships_drawn": len(rels) - len(rel_errors),
        "warnings": rel_errors,
    }


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #
def parse_args():
    p = argparse.ArgumentParser(description="Render a CSDM v5 PowerPoint slide from a JSON spec.")
    p.add_argument("spec_pos",  nargs="?", help="Positional spec path")
    p.add_argument("out_pos",   nargs="?", help="Positional output path")
    p.add_argument("--spec")
    p.add_argument("--out")
    p.add_argument("--catalog", default=str(DEFAULT_CATALOG))
    p.add_argument("--dump-catalog", action="store_true",
                   help="List all known canonical objects and roles, then exit.")
    return p.parse_args()


def main():
    args = parse_args()
    catalog = load_catalog(Path(args.catalog))
    if args.dump_catalog:
        print("== DOMAINS ==")
        for d in catalog["domains"]:
            print(f"  {d['id']:30s} fill=#{d['fill']}  label={d['label']}")
        print("\n== OBJECTS ==")
        for o in catalog["objects"]:
            syn = ", ".join(o.get("synonyms", []))
            print(f"  {o['canonical']:42s}  domain={o['domain']:25s}  synonyms=[{syn}]")
        print("\n== ROLES ==")
        for r in catalog["roles"]:
            print(f"  {r['canonical']}")
        print("\n== STANDARD RELATIONSHIPS ==")
        for rel in catalog.get("standard_relationships", []):
            print(f"  {rel['from']} -> {rel['to']} : {rel['label']}")
        return 0

    spec_path = Path(args.spec or args.spec_pos or "")
    out_path  = Path(args.out  or args.out_pos  or "")
    if not spec_path.exists():
        sys.stderr.write(f"Spec not found: {spec_path}\n")
        return 2
    if not out_path:
        sys.stderr.write("Output path required\n")
        return 2
    with open(spec_path) as f:
        spec = json.load(f)
    try:
        result = render(spec, catalog, out_path)
    except ValueError as e:
        sys.stderr.write(f"{e}\n")
        return 2
    print(json.dumps(result, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
