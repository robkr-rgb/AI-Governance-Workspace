"""Van Oord deck helpers.

Import:
    from vanoord_helpers import (
        BRAND, load_template, add_title, text_box, grid,
        bottom_strip, save_and_render
    )

Use `load_template(output_path)` as the first call in every deck builder,
then compose slides with `add_title`, `text_box`, `grid`, and `bottom_strip`.
"""
from __future__ import annotations
import os
import shutil
import subprocess
from dataclasses import dataclass
from typing import Any, Iterable

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Brand constants, extracted directly from the Van Oord template theme.
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class Brand:
    navy: RGBColor = RGBColor(0x25, 0x3A, 0x79)
    navy_mid: RGBColor = RGBColor(0x1C, 0x4B, 0x8A)
    navy_light: RGBColor = RGBColor(0x24, 0x60, 0xA7)
    orange: RGBColor = RGBColor(0xFF, 0x87, 0x02)
    amber: RGBColor = RGBColor(0xFF, 0xA4, 0x3F)
    steel: RGBColor = RGBColor(0xA8, 0xB0, 0xC9)
    light: RGBColor = RGBColor(0xF2, 0xF2, 0xF2)
    grey_mid: RGBColor = RGBColor(0xD9, 0xDD, 0xE3)
    dark: RGBColor = RGBColor(0x26, 0x26, 0x26)
    white: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)
    muted: RGBColor = RGBColor(0x55, 0x55, 0x55)

    major_font: str = "Arial Black"
    minor_font: str = "Arial"
    mono_font: str = "Consolas"


BRAND = Brand()


# Standard slide geometry (16:9, 13.33 x 7.5 inches).
MARGIN_L = Emu(594360)
MARGIN_R = Emu(594360)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


# ---------------------------------------------------------------------------
# Template loading.
# ---------------------------------------------------------------------------

def _template_path() -> str:
    """Return absolute path to the bundled Van Oord template."""
    here = os.path.dirname(os.path.abspath(__file__))
    candidates = [
        os.path.join(here, "..", "assets", "vanoord_template.pptx"),
        os.path.join(here, "assets", "vanoord_template.pptx"),
    ]
    for c in candidates:
        if os.path.isfile(c):
            return os.path.abspath(c)
    raise FileNotFoundError(
        "vanoord_template.pptx not found. Expected next to this script in assets/."
    )


def load_template(output_path: str) -> Presentation:
    """Copy the bundled template to `output_path` and return a Presentation.

    The returned Presentation already has Van Oord masters, layouts, theme
    colors, fonts and logo in place. Add content slides by calling
    `prs.slides.add_slide(prs.slide_layouts[i])`.
    """
    src = _template_path()
    os.makedirs(os.path.dirname(os.path.abspath(output_path)) or ".", exist_ok=True)
    shutil.copyfile(src, output_path)
    return Presentation(output_path)


# ---------------------------------------------------------------------------
# Text box primitive.
# ---------------------------------------------------------------------------

def text_box(
    slide,
    left, top, w, h,
    text,
    *,
    size: float = 12,
    bold: bool = False,
    color: RGBColor = BRAND.dark,
    font: str = None,
    fill: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    border_color: RGBColor | None = None,
    pad: float = 0.08,
):
    """Add a rectangle with a text frame. The universal building block.

    `text` accepts either a plain string or a list. List items may be:
      - plain string -> inherits the outer styling
      - (string, opts_dict) -> run-level overrides
           opts keys: size, bold, color, font
      - separate paragraph when the list item is a top-level string or tuple

    Returns the shape so the caller can tweak further.
    """
    font = font or BRAND.minor_font

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.line.fill.background()
    if border_color is not None:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(pad)
    tf.margin_right = Inches(pad)
    tf.margin_top = Inches(pad * 0.6)
    tf.margin_bottom = Inches(pad * 0.6)
    tf.vertical_anchor = anchor

    if isinstance(text, list):
        for i, item in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if isinstance(item, tuple):
                t, opts = item
                run = p.add_run()
                run.text = t
                run.font.name = opts.get("font", font)
                run.font.size = Pt(opts.get("size", size))
                run.font.bold = opts.get("bold", bold)
                run.font.color.rgb = opts.get("color", color)
            else:
                run = p.add_run()
                run.text = item
                run.font.name = font
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = color
    else:
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    return shape


# ---------------------------------------------------------------------------
# Title block.
# ---------------------------------------------------------------------------

def add_title(slide, title: str, subtitle: str | None = None):
    """Place the Van Oord brand title and subtitle in the standard position."""
    text_box(slide,
             Emu(714977), Emu(373978), Emu(10695889), Emu(340000),
             title,
             size=24, bold=True, color=BRAND.navy, font=BRAND.minor_font)
    if subtitle:
        text_box(slide,
                 Emu(714977), Emu(720000), Emu(10697600), Emu(420000),
                 subtitle,
                 size=13, color=BRAND.muted, font=BRAND.minor_font)


# ---------------------------------------------------------------------------
# Grid (custom table).
# ---------------------------------------------------------------------------

def grid(
    slide,
    left, top, w, h,
    headers: list,
    rows: list[list],
    *,
    col_weights: list[float] | None = None,
    header_fill: RGBColor = BRAND.navy,
    header_color: RGBColor = BRAND.white,
    header_size: float = 10,
    cell_size: float = 9,
    zebra: bool = True,
    header_h_in: float = 0.38,
):
    """Render a dense grid using shapes rather than PPTX native tables.

    `headers`: list of strings.
    `rows`: list of rows. Each row is a list where each cell is either:
        - a string                              -> default cell style
        - a tuple (string, opts_dict) -> per-cell overrides:
             opts keys: size, bold, color, mono (bool), font
    `col_weights`: optional list of floats controlling relative widths.
    """
    n_cols = len(headers)
    weights = col_weights or [1] * n_cols
    total = sum(weights)
    col_widths = [Emu(int(w * cw / total)) for cw in weights]
    x_offsets = []
    x = left
    for cw in col_widths:
        x_offsets.append(x)
        x += cw

    hdr_h = Emu(int(Inches(header_h_in)))
    row_h_val = (h - hdr_h) / max(len(rows), 1)
    row_h = Emu(int(row_h_val))

    # Header row
    for i, hdr in enumerate(headers):
        text_box(slide, x_offsets[i], top, col_widths[i], hdr_h,
                 hdr,
                 size=header_size, bold=True, color=header_color,
                 fill=header_fill,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT, pad=0.08)

    # Data rows
    for ri, row in enumerate(rows):
        ytop = top + hdr_h + Emu(int(row_h_val * ri))
        row_fill = BRAND.light if (zebra and ri % 2 == 0) else BRAND.white
        for ci, cell in enumerate(row):
            if isinstance(cell, tuple):
                txt, opts = cell
            else:
                txt, opts = cell, {}
            mono = opts.get("mono", False)
            text_box(slide, x_offsets[ci], ytop, col_widths[ci], row_h,
                     txt,
                     size=opts.get("size", cell_size),
                     bold=opts.get("bold", False),
                     color=opts.get("color", BRAND.dark),
                     fill=row_fill,
                     font=opts.get("font", BRAND.mono_font if mono else BRAND.minor_font),
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT, pad=0.08)


# ---------------------------------------------------------------------------
# Bottom strip.
# ---------------------------------------------------------------------------

def bottom_strip(slide, text: str, color: str = "navy", *, bold: bool = True):
    """Add a one-line key-message strip across the bottom of the slide.

    `color` is 'navy' (default) or 'orange'.
    """
    fill = BRAND.navy_mid if color == "navy" else BRAND.orange
    text_box(slide, MARGIN_L, Emu(6400000), CONTENT_W, Emu(300000),
             text,
             size=11, bold=bold, color=BRAND.white, fill=fill,
             anchor=MSO_ANCHOR.MIDDLE, pad=0.1,
             align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Save and render for QA.
# ---------------------------------------------------------------------------

def save_and_render(prs: Presentation, pptx_path: str, *, render_pdf: bool = True) -> str:
    """Save the Presentation and optionally render to PDF via LibreOffice.

    Returns the PDF path, or '' if rendering was skipped or failed.
    """
    prs.save(pptx_path)
    if not render_pdf:
        return ""

    out_dir = os.path.dirname(os.path.abspath(pptx_path)) or "."
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             "--outdir", out_dir, pptx_path],
            check=True, capture_output=True, timeout=120,
        )
    except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
        return ""
    pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
    return pdf_path if os.path.isfile(pdf_path) else ""
