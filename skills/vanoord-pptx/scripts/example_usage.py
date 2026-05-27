"""Minimal example showing how to use the Van Oord helpers.

Run:
    python example_usage.py
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from vanoord_helpers import (
    BRAND, load_template, add_title, text_box, grid, bottom_strip,
    save_and_render, MARGIN_L, CONTENT_W,
)
from pptx.util import Emu, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

prs = load_template("example_output.pptx")

# Cover slide using Title Slide layout (index 1)
cover = prs.slides.add_slide(prs.slide_layouts[1])
# Title Slide layout placeholders: fill the first text placeholder if it exists
for ph in cover.placeholders:
    if "Title" in ph.name and ph.has_text_frame:
        ph.text_frame.text = "Example Van Oord Deck"
        break

# Content slide on CUSTOM layout (index 12)
s = prs.slides.add_slide(prs.slide_layouts[12])
add_title(s,
          "Example dense content slide",
          "Table, key message, and one-line conclusion")

headers = ["Item", "Owner", "Status", "Identifier"]
rows = [
    ["Alpha", "Operations",   "On track",  ("ID-0001", {"mono": True})],
    ["Beta",  "Engineering",  "Delayed",   ("ID-0002", {"mono": True})],
    ["Gamma", "Finance",      "Complete",  ("ID-0003", {"mono": True})],
]
grid(s, MARGIN_L, Emu(1400000), CONTENT_W, Emu(4700000),
     headers, rows, col_weights=[1.2, 1.8, 1.2, 1.6])

bottom_strip(s, "Structure first, content second. This is how we stay on brand.", color="orange")

pdf = save_and_render(prs, "example_output.pptx")
print("Saved PPTX, PDF at:", pdf or "(PDF skipped, soffice missing)")
